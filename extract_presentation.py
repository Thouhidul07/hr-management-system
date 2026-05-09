#!/usr/bin/env python3
"""
Script to extract content from PowerPoint presentation
"""

from pptx import Presentation
import os

def extract_presentation_content(pptx_path):
    """Extract all content from a PowerPoint file"""
    
    if not os.path.exists(pptx_path):
        print(f"Error: File not found at {pptx_path}")
        return
    
    try:
        # Load the presentation
        prs = Presentation(pptx_path)
        
        print("=" * 80)
        print(f"PRESENTATION EXTRACTED: {os.path.basename(pptx_path)}")
        print("=" * 80)
        print()
        
        # Iterate through slides
        for slide_num, slide in enumerate(prs.slides, 1):
            print(f"\n{'=' * 80}")
            print(f"SLIDE {slide_num}")
            print(f"{'=' * 80}\n")
            
            # Extract text from shapes
            for shape_num, shape in enumerate(slide.shapes):
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:  # Only print non-empty text
                        print(f"{text}\n")
                
                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    print("TABLE CONTENT:")
                    for row in table.rows:
                        row_data = [cell.text for cell in row.cells]
                        print(" | ".join(row_data))
                    print()
                
                # Extract text from text frames
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            print(para_text)
            
        print(f"\n{'=' * 80}")
        print(f"Total Slides: {len(prs.slides)}")
        print(f"{'=' * 80}\n")
        
    except Exception as e:
        print(f"Error reading presentation: {e}")

if __name__ == "__main__":
    pptx_file = "HR_Management_Presentation.pptx"
    extract_presentation_content(pptx_file)
